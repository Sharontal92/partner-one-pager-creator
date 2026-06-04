#!/usr/bin/env python3
"""
Partner One Pager — PPTX Generator
Produces a 2-slide PowerPoint following the Optimove partner one-pager template.

Usage:
    python generate_pptx.py <partner_data.json> <output_dir> [--aspect 16:9|4:3]
"""

import sys
import json
import os
import re
import argparse
import pathlib
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Mm
    import pptx.oxml.ns as pns
    from lxml import etree
except ImportError:
    print("Installing python-pptx...")
    os.system("pip install python-pptx lxml --break-system-packages -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Mm
    import pptx.oxml.ns as pns
    from lxml import etree

# ── Brand Colors ───────────────────────────────────────────────────────────────
MIDNIGHT_VIOLET = RGBColor(0x30, 0x2c, 0x69)
LIME_GLOW       = RGBColor(0xDF, 0xF6, 0x70)
SHADOW_BLACK    = RGBColor(0x11, 0x11, 0x11)
CLOUD_TINT      = RGBColor(0xEF, 0xEF, 0xEF)
STORM_TINT      = RGBColor(0xD3, 0xD3, 0xD3)
SIGNAL_MIST     = RGBColor(0x9E, 0x97, 0xCB)
PURE_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

POPPINS = "Poppins"

SKILL_DIR = Path(__file__).parent.parent
# Try several candidate paths for the font directory
_font_candidates = [
    SKILL_DIR.parent / "optimove-brand-bible" / "assets" / "fonts" / "Poppins",
    SKILL_DIR.parent.parent / "optimove-brand-bible" / "assets" / "fonts" / "Poppins",
    pathlib.Path("/sessions").glob("*/mnt/.claude/skills/optimove-brand-bible/assets/fonts/Poppins"),
]
FONT_DIR = next((p for p in _font_candidates if isinstance(p, pathlib.Path) and p.exists()), _font_candidates[0])

# ── Slide dimensions ───────────────────────────────────────────────────────────
ASPECT_RATIOS = {
    "16:9": (Cm(33.867), Cm(19.05)),   # 13.33" × 7.5"
    "4:3":  (Cm(25.4),   Cm(19.05)),   # 10"    × 7.5"
}

MARGIN = Cm(1.9)   # ~0.75 inch


# ── Helpers ────────────────────────────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text,
                font_name=POPPINS, size_pt=12, bold=False, italic=False,
                color=SHADOW_BLACK, align=PP_ALIGN.LEFT,
                wrap=True, auto_fit=False):
    """Add a text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(size_pt)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return txBox


def add_colored_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()  # No line
    return shape


def add_circle(slide, left, top, diameter, fill_color):
    """Add a filled ellipse (used as icon background)."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        9,  # Oval
        left, top, diameter, diameter
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_divider(slide, left, top, width, color=STORM_TINT):
    """Add a thin horizontal line."""
    line = slide.shapes.add_shape(1, left, top, width, Pt(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_multiline_textbox(slide, left, top, width, height,
                           lines: list,
                           font_name=POPPINS, size_pt=12, bold=False,
                           italic=False, color=SHADOW_BLACK,
                           align=PP_ALIGN.LEFT, line_spacing_pt=None):
    """Add a text box with multiple paragraphs."""
    from pptx.util import Pt as _Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    for i, (text, kwargs) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = kwargs.get("align", align)
        run = p.add_run()
        run.text            = text
        run.font.name       = font_name
        run.font.size       = _Pt(kwargs.get("size_pt", size_pt))
        run.font.bold       = kwargs.get("bold", bold)
        run.font.italic     = kwargs.get("italic", italic)
        run.font.color.rgb  = kwargs.get("color", color)

        if line_spacing_pt:
            from pptx.util import Pt as _Pt2
            p.line_spacing = _Pt2(line_spacing_pt)

    return txBox


# ── Slide builders ─────────────────────────────────────────────────────────────
def build_slide1(prs, data, slide_w, slide_h):
    """Slide 1: Header + Headline + Intro + Capabilities."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PURE_WHITE

    content_w = slide_w - 2 * MARGIN
    y = MARGIN

    # ── Header ────────────────────────────────────────────────────────────────
    header_h = Cm(1.4)
    # Optimove logo text (left)
    add_textbox(slide,
                left=MARGIN, top=y, width=content_w / 2, height=header_h,
                text="OPTIMOVE",
                size_pt=14, bold=True, color=MIDNIGHT_VIOLET, align=PP_ALIGN.LEFT)
    # Partner name (right)
    add_textbox(slide,
                left=MARGIN + content_w / 2, top=y,
                width=content_w / 2, height=header_h,
                text=data["partner_name"],
                size_pt=12, bold=True, color=SHADOW_BLACK, align=PP_ALIGN.RIGHT)
    # Divider under header
    y += header_h
    add_divider(slide, MARGIN, y, content_w)
    y += Cm(0.8)

    # ── Headline ──────────────────────────────────────────────────────────────
    headline_h = Cm(3.5)
    add_textbox(slide,
                left=MARGIN, top=y, width=content_w, height=headline_h,
                text=data["headline"],
                size_pt=34, bold=True, color=SHADOW_BLACK, align=PP_ALIGN.LEFT)
    y += headline_h

    # ── Intro ─────────────────────────────────────────────────────────────────
    intro = data.get("intro", "")
    intro_paras = [p.strip() for p in intro.split("\n\n") if p.strip()]
    if not intro_paras:
        intro_paras = [intro]

    intro_lines = []
    for i, para in enumerate(intro_paras):
        intro_lines.append((para, {"size_pt": 11}))
        if i < len(intro_paras) - 1:
            intro_lines.append(("", {"size_pt": 6}))

    intro_h = Cm(2.8)
    add_multiline_textbox(slide,
                          left=MARGIN, top=y, width=content_w, height=intro_h,
                          lines=intro_lines,
                          size_pt=11, color=SHADOW_BLACK)
    y += intro_h + Cm(0.4)

    # Divider
    add_divider(slide, MARGIN, y, content_w)
    y += Cm(0.5)

    # ── Capabilities ─────────────────────────────────────────────────────────
    bg_colors = [MIDNIGHT_VIOLET, LIME_GLOW, SIGNAL_MIST]
    sym_colors = [PURE_WHITE, MIDNIGHT_VIOLET, PURE_WHITE]

    icon_d   = Cm(2.2)
    icon_col = Cm(3.0)
    cap_h    = Cm(2.5)
    gap      = Cm(0.4)

    for i, cap in enumerate(data.get("capabilities", [])[:3]):
        # Colored circle
        circle_left = MARGIN + (icon_col - icon_d) / 2
        circle_top  = y + (cap_h - icon_d) / 2
        add_circle(slide, circle_left, circle_top, icon_d, bg_colors[i % len(bg_colors)])

        # Title + description
        text_left = MARGIN + icon_col
        text_w    = content_w - icon_col

        lines = [
            (cap["title"],       {"bold": True,  "size_pt": 13, "color": MIDNIGHT_VIOLET}),
            ("",                 {"size_pt": 4}),
            (cap["description"], {"bold": False, "size_pt": 10, "color": SHADOW_BLACK}),
        ]
        add_multiline_textbox(slide,
                              left=text_left, top=y,
                              width=text_w, height=cap_h,
                              lines=lines)
        y += cap_h + gap

    return slide


def build_slide2(prs, data, slide_w, slide_h):
    """Slide 2: Use cases + Social proof + KPIs."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PURE_WHITE

    content_w = slide_w - 2 * MARGIN
    y = MARGIN

    # ── Header ────────────────────────────────────────────────────────────────
    header_h = Cm(1.4)
    add_textbox(slide,
                left=MARGIN, top=y, width=content_w / 2, height=header_h,
                text="OPTIMOVE",
                size_pt=14, bold=True, color=MIDNIGHT_VIOLET, align=PP_ALIGN.LEFT)
    add_textbox(slide,
                left=MARGIN + content_w / 2, top=y,
                width=content_w / 2, height=header_h,
                text=data["partner_name"],
                size_pt=12, bold=True, color=SHADOW_BLACK, align=PP_ALIGN.RIGHT)
    y += header_h
    add_divider(slide, MARGIN, y, content_w)
    y += Cm(0.6)

    # ── Use Cases section ─────────────────────────────────────────────────────
    add_textbox(slide,
                left=MARGIN, top=y, width=content_w, height=Cm(1.0),
                text="Seamlessly power any use case",
                size_pt=18, bold=True, color=SHADOW_BLACK, align=PP_ALIGN.LEFT)
    y += Cm(1.1)

    uc_h    = Cm(1.3)
    uc_gap  = Cm(0.2)
    for uc in data.get("use_cases", [])[:4]:
        lines = [
            (uc["title"],       {"bold": True,  "size_pt": 11, "color": MIDNIGHT_VIOLET}),
            (uc["description"], {"bold": False, "size_pt": 10, "color": SHADOW_BLACK}),
        ]
        add_multiline_textbox(slide,
                              left=MARGIN, top=y,
                              width=content_w, height=uc_h,
                              lines=lines)
        y += uc_h + uc_gap

    y += Cm(0.3)
    add_divider(slide, MARGIN, y, content_w)
    y += Cm(0.4)

    # ── Social Proof ──────────────────────────────────────────────────────────
    trusted_by = data.get("trusted_by", "1,200+")
    trusted_h = Cm(0.8)
    add_textbox(slide,
                left=MARGIN, top=y, width=content_w, height=trusted_h,
                text=f"Trusted by {trusted_by} Brands",
                size_pt=16, bold=True, color=SHADOW_BLACK, align=PP_ALIGN.CENTER)
    y += trusted_h + Cm(0.3)

    clients = (data.get("client_logos", []) + [""] * 6)[:6]
    cell_w  = content_w / 3
    cell_h  = Cm(1.0)
    for row_i in range(2):
        for col_i in range(3):
            name = clients[row_i * 3 + col_i]
            cx = MARGIN + col_i * cell_w
            cy = y + row_i * cell_h
            # Cell border
            border_rect = add_colored_rect(slide, cx, cy, cell_w, cell_h,
                                           PURE_WHITE, line_color=STORM_TINT)
            if name:
                add_textbox(slide,
                            left=cx, top=cy, width=cell_w, height=cell_h,
                            text=name,
                            size_pt=9, bold=True, color=SHADOW_BLACK,
                            align=PP_ALIGN.CENTER)
    y += cell_h * 2 + Cm(0.3)

    # Optional quote
    quote = data.get("client_quote")
    if quote and quote.get("text") and y + Cm(1.5) < slide_h - MARGIN - Cm(3):
        add_divider(slide, MARGIN, y, content_w)
        y += Cm(0.3)
        q_h = Cm(1.2)
        quote_text = f"“{quote['text']}”"
        lines = [(quote_text, {"italic": True, "color": MIDNIGHT_VIOLET, "size_pt": 10})]
        if quote.get("attribution"):
            lines += [("", {"size_pt": 3}),
                      (quote["attribution"], {"bold": True, "color": SHADOW_BLACK, "size_pt": 8})]
        add_multiline_textbox(slide,
                              left=MARGIN + Cm(1), top=y,
                              width=content_w - Cm(2), height=q_h,
                              lines=lines, align=PP_ALIGN.CENTER)
        y += q_h + Cm(0.2)

    add_divider(slide, MARGIN, y, content_w)
    y += Cm(0.4)

    # ── KPIs ──────────────────────────────────────────────────────────────────
    cta = data.get("cta", "Do more when you go Positionless")
    add_textbox(slide,
                left=MARGIN, top=y, width=content_w, height=Cm(0.9),
                text=cta,
                size_pt=16, bold=True, color=SHADOW_BLACK, align=PP_ALIGN.CENTER)
    y += Cm(1.0)

    kpis  = data.get("kpis", [])[:3]
    kpi_h = slide_h - MARGIN - y - Cm(0.2)
    if kpis and kpi_h > Cm(0.5):
        n     = len(kpis)
        kw    = content_w / n
        for i, kpi in enumerate(kpis):
            kx = MARGIN + i * kw
            lines = [
                (kpi["value"], {"bold": True, "italic": True, "size_pt": 28,
                                "color": MIDNIGHT_VIOLET, "align": PP_ALIGN.CENTER}),
                ("",           {"size_pt": 4}),
                (kpi["label"], {"bold": False, "size_pt": 9,
                                "color": SHADOW_BLACK, "align": PP_ALIGN.CENTER}),
            ]
            add_multiline_textbox(slide,
                                  left=kx, top=y, width=kw, height=kpi_h,
                                  lines=lines, align=PP_ALIGN.CENTER)
            # Vertical divider between KPIs
            if i < n - 1:
                div_x = kx + kw
                add_colored_rect(slide, div_x, y, Pt(0.5), kpi_h, STORM_TINT)

    return slide


# ── Main ───────────────────────────────────────────────────────────────────────
def generate(data_path: str, output_dir: str, aspect: str = "16:9") -> str:
    with open(data_path) as f:
        data = json.load(f)

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if aspect not in ASPECT_RATIOS:
        print(f"⚠️  Unknown aspect ratio '{aspect}', defaulting to 16:9")
        aspect = "16:9"

    slide_w, slide_h = ASPECT_RATIOS[aspect]

    partner  = data["partner_name"]
    safe     = re.sub(r'[^\w\s-]', '', partner).strip()
    out_file = output_dir / f"Optimove x {safe} - One Pager.pptx"

    prs = Presentation()
    prs.slide_width  = slide_w
    prs.slide_height = slide_h

    build_slide1(prs, data, slide_w, slide_h)
    build_slide2(prs, data, slide_w, slide_h)

    prs.save(str(out_file))
    print(f"✅ PPTX generated ({aspect}): {out_file}")
    return str(out_file)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("data",       help="Path to partner_data.json")
    parser.add_argument("output_dir", help="Output directory")
    parser.add_argument("--aspect",   default="16:9",
                        choices=["16:9", "4:3"],
                        help="Slide aspect ratio")
    args = parser.parse_args()
    generate(args.data, args.output_dir, args.aspect)
